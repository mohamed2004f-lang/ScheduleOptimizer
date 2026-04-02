import sys, os
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
from models.models import Course
from flask import Blueprint, request, jsonify, Response, current_app, session, send_file
from backend.core.auth import login_required, role_required
from collections import defaultdict
from .utilities import get_connection, table_to_dicts, df_from_query, excel_response_from_df, pdf_response_from_html
import io
import base64
from datetime import datetime

courses_bp = Blueprint("courses", __name__)


def _is_instructor_or_supervisor_view_only() -> bool:
    role = (session.get("user_role") or "").strip()
    return role == "supervisor" or (role == "instructor") or (role == "instructor" and int(session.get("is_supervisor") or 0) == 1)

@courses_bp.route("/list")
@login_required
def list_courses():
    # يرجع جدول courses، وإذا غير موجود يرجع من schedule
    with get_connection() as conn:
        cur = conn.cursor()
        try:
            try:
                cols = [r[1] for r in cur.execute("PRAGMA table_info(courses)").fetchall()]
            except Exception:
                cols = []
            has_cat = "category" in cols
            has_archived = "is_archived" in cols
            sel = "SELECT DISTINCT course_name, course_code, units"
            if has_cat:
                sel += ", COALESCE(category,'required') AS category"
            if has_archived:
                sel += ", COALESCE(is_archived,0) AS is_archived"
            sel += " FROM courses WHERE COALESCE(course_name,'') <> '' ORDER BY course_name"
            rows = cur.execute(sel).fetchall()
            # إزالة التكرار البرمجيًا أيضًا (احتياطي)
            seen = set()
            courses = []
            for r in rows:
                cname = r[0]
                key = cname.strip().lower() if cname else ""
                if not cname or key in seen:
                    continue
                seen.add(key)
                c = Course(r[0], r[1], r[2])
                try:
                    setattr(c, "category", (r[3] if has_cat else "required") or "required")
                except Exception:
                    setattr(c, "category", "required")
                try:
                    archived_idx = 4 if has_cat else 3
                    setattr(c, "is_archived", int(r[archived_idx] or 0) if has_archived else 0)
                except Exception:
                    setattr(c, "is_archived", 0)
                courses.append(c)
        except Exception:
            rows = cur.execute("SELECT DISTINCT course_name FROM schedule WHERE COALESCE(course_name,'') <> '' ORDER BY course_name").fetchall()
            seen = set()
            courses = []
            for r in rows:
                cname = r[0]
                key = cname.strip().lower() if cname else ""
                if not cname or key in seen:
                    continue
                seen.add(key)
                c = Course(r[0], "", 0)
                setattr(c, "category", "required")
                courses.append(c)
    return jsonify([c.__dict__ for c in courses])

@courses_bp.route("/add", methods=["POST"])
@role_required("admin", "admin_main", "head_of_department")
def add_course():
    data = request.get_json(force=True)
    cname = (data.get("course_name") or "").strip()
    code = (data.get("course_code") or "").strip()
    try:
        units = int(data.get("units", 0) or 0)
    except (TypeError, ValueError):
        units = 0
    category = (data.get("category") or "required").strip() or "required"
    if category not in ("required", "elective_major", "elective_free"):
        category = "required"
    if not cname:
        return jsonify({"status": "error", "message": "اسم المقرر (course_name) مطلوب"}), 400
    with get_connection() as conn:
        cur = conn.cursor()
        cur.execute("""
            CREATE TABLE IF NOT EXISTS courses (
                course_name TEXT PRIMARY KEY,
                course_code TEXT,
                units INTEGER
            )
        """)
        # منع تكرار الاسم (تطبيع بسيط lower/strip)
        row = cur.execute(
            "SELECT course_name FROM courses WHERE LOWER(TRIM(course_name)) = LOWER(TRIM(?))",
            (cname,),
        ).fetchone()
        if row:
            return jsonify({"status": "error", "message": "يوجد مقرر آخر بنفس الاسم. استخدم زر \"تحرير\" لتعديله."}), 400

        # منع تكرار الرمز إذا تم إدخاله
        if code:
            row = cur.execute(
                "SELECT course_name FROM courses WHERE COALESCE(course_code,'') <> '' AND LOWER(TRIM(course_code)) = LOWER(TRIM(?))",
                (code,),
            ).fetchone()
            if row:
                return jsonify(
                    {
                        "status": "error",
                        "message": f"يوجد مقرر آخر بنفس الرمز ({row['course_name']}). الرجاء اختيار رمز مختلف.",
                    }
                ), 400

        # تأكد من وجود عمود category (قواعد قديمة)
        try:
            cols = [r[1] for r in cur.execute("PRAGMA table_info(courses)").fetchall()]
        except Exception:
            cols = []
        if "category" in cols:
            cur.execute(
                "INSERT INTO courses (course_name, course_code, units, category) VALUES (?, ?, ?, ?)",
                (cname, code, units, category),
            )
        else:
            cur.execute(
                "INSERT INTO courses (course_name, course_code, units) VALUES (?, ?, ?)",
                (cname, code, units),
            )
        conn.commit()
    return jsonify({"status": "ok", "message": "تم إضافة المقرر"}), 200

@courses_bp.route("/update", methods=["POST"])
@role_required("admin", "admin_main", "head_of_department")
def update_course():
    data = request.get_json(force=True)
    old_name = (data.get("old_course_name") or "").strip()
    new_name = (data.get("new_course_name") or "").strip()
    new_units = data.get("units")
    new_code = (data.get("course_code") or "").strip()
    category = (data.get("category") or "").strip()
    if not old_name or not new_name:
        return jsonify({"status": "error", "message": "old_course_name و new_course_name مطلوبة"}), 400

    with get_connection() as conn:
        cur = conn.cursor()
        # منع تكرار الاسم الجديد (باستثناء نفس المقرر)
        row = cur.execute(
            "SELECT course_name FROM courses WHERE LOWER(TRIM(course_name)) = LOWER(TRIM(?)) AND course_name <> ?",
            (new_name, old_name),
        ).fetchone()
        if row:
            return jsonify({"status": "error", "message": "لا يمكن تغيير الاسم لأنه مستخدم لمقرر آخر."}), 400

        # منع تكرار الرمز الجديد إن وجد
        if new_code:
            row = cur.execute(
                """
                SELECT course_name FROM courses
                WHERE COALESCE(course_code,'') <> ''
                  AND LOWER(TRIM(course_code)) = LOWER(TRIM(?))
                  AND course_name <> ?
                """,
                (new_code, old_name),
            ).fetchone()
            if row:
                return jsonify(
                    {
                        "status": "error",
                        "message": f"الرمز مستخدم لمقرر آخر ({row['course_name']}). اختر رمزاً مختلفاً.",
                    }
                ), 400

        # تأكد من وجود عمود category (قواعد قديمة)
        try:
            cols = [r[1] for r in cur.execute("PRAGMA table_info(courses)").fetchall()]
        except Exception:
            cols = []
        has_cat = "category" in cols
        cat_value = category if category in ("required", "elective_major", "elective_free") else None
        if has_cat:
            if cat_value is None:
                cur.execute(
                    "UPDATE courses SET course_name=?, course_code=?, units=? WHERE course_name=?",
                    (new_name, new_code or "", (int(new_units) if new_units is not None else None), old_name),
                )
            else:
                cur.execute(
                    "UPDATE courses SET course_name=?, course_code=?, units=?, category=? WHERE course_name=?",
                    (new_name, new_code or "", (int(new_units) if new_units is not None else None), cat_value, old_name),
                )
        else:
            cur.execute(
                "UPDATE courses SET course_name=?, course_code=?, units=? WHERE course_name=?",
                (new_name, new_code or "", (int(new_units) if new_units is not None else None), old_name),
            )

        # تحديث جميع الجداول التي تعتمد على اسم المقرر
        for tbl in ("grades", "schedule", "registrations", "enrollment_plan_items", "exams"):
            try:
                cur.execute(f"UPDATE {tbl} SET course_name=? WHERE course_name=?", (new_name, old_name))
            except Exception:
                pass

        cur.execute("UPDATE prereqs SET course_name=? WHERE course_name=?", (new_name, old_name))
        cur.execute("UPDATE prereqs SET required_course_name=? WHERE required_course_name=?", (new_name, old_name))

        if new_units is not None or new_code is not None:
            try:
                if new_units is not None:
                    cur.execute("UPDATE grades SET units=? WHERE course_name=?", (int(new_units), new_name))
                if new_code is not None:
                    cur.execute(
                        "UPDATE grades SET course_code=? WHERE course_name=?",
                        (new_code or "", new_name),
                    )
            except Exception:
                pass

        # أي تعديل على المقررات (الاسم/الرمز/الوحدات) يجعل نتائج التحسين الحالية قديمة
        try:
            cur.execute("DELETE FROM optimized_schedule")
        except Exception:
            pass
        try:
            cur.execute("DELETE FROM conflict_report")
        except Exception:
            pass

        conn.commit()
    return jsonify({"status": "ok", "message": "تم تعديل بيانات المقرر"}), 200

@courses_bp.route("/delete", methods=["POST"])
@role_required("admin", "admin_main", "head_of_department")
def delete_course():
    data = request.get_json(force=True)
    cname = data.get("course_name")
    if not cname:
        return jsonify({"status": "error", "message": "course_name مطلوب"}), 400
    with get_connection() as conn:
        cur = conn.cursor()
        # تحقق ارتباطات أكاديمية - لا نحذف صلباً عند وجود آثار
        links = {}
        for tbl in ("grades", "registrations", "schedule", "enrollment_plan_items", "exams", "prereqs"):
            try:
                if tbl == "prereqs":
                    row = cur.execute(
                        "SELECT COUNT(*) FROM prereqs WHERE course_name = ? OR required_course_name = ?",
                        (cname, cname),
                    ).fetchone()
                else:
                    row = cur.execute(
                        f"SELECT COUNT(*) FROM {tbl} WHERE course_name = ?",
                        (cname,),
                    ).fetchone()
                links[tbl] = int(row[0] or 0) if row else 0
            except Exception:
                links[tbl] = 0

        has_links = any(v > 0 for v in links.values())
        # ensure archive column exists
        try:
            cols = [r[1] for r in cur.execute("PRAGMA table_info(courses)").fetchall()]
        except Exception:
            cols = []
        if "is_archived" not in cols:
            try:
                cur.execute("ALTER TABLE courses ADD COLUMN is_archived INTEGER NOT NULL DEFAULT 0")
            except Exception:
                pass

        if has_links:
            cur.execute("UPDATE courses SET is_archived = 1 WHERE course_name = ?", (cname,))
            conn.commit()
            return jsonify({
                "status": "ok",
                "archived": True,
                "message": "تمت أرشفة المقرر بدلاً من الحذف لأنه مرتبط ببيانات أكاديمية تاريخية.",
                "links": links,
            }), 200

        # حذف صلب فقط إذا لا يوجد أي ارتباط
        cur.execute("DELETE FROM courses WHERE course_name = ?", (cname,))
        cur.execute("DELETE FROM prereqs WHERE course_name = ? OR required_course_name = ?", (cname, cname))
        try:
            cur.execute("DELETE FROM optimized_schedule")
        except Exception:
            pass
        try:
            cur.execute("DELETE FROM conflict_report")
        except Exception:
            pass
        conn.commit()
    return jsonify({"status": "ok", "archived": False, "message": "تم حذف المقرر (لا توجد له ارتباطات)."}), 200

# المتطلبات (Prereqs) - يدعم زوج واحد أو دفعة items[]
@courses_bp.route("/prereqs/add", methods=["POST"])
@role_required("admin", "admin_main", "head_of_department")
def add_prereq():
    """
    Accepts:
    - single object: {"course_name":"A","required_course_name":"B"}
    - or batch: {"items":[{"course_name":"A","required_course_name":"B"}, ...]}

    Response contains lists: added, ignored (duplicates), missing, errors
    """
    data = request.get_json(force=True) or {}

    # normalize incoming items into a list of pairs
    items = []
    if isinstance(data, dict) and "items" in data and isinstance(data["items"], list):
        for it in data["items"]:
            c = (it.get("course_name") or "").strip()
            r = (it.get("required_course_name") or "").strip()
            if c and r:
                items.append((c, r))
    else:
        # allow single pair payload
        c = (data.get("course_name") or "").strip()
        r = (data.get("required_course_name") or "").strip()
        if c and r:
            items.append((c, r))

    if not items:
        return jsonify({"status":"error","message":"يرجى تمرير course_name و required_course_name أو مصفوفة items"}), 400

    added = []
    ignored = []
    missing = []
    errors = []

    with get_connection() as conn:
        cur = conn.cursor()
        # Ensure prereqs table exists
        cur.execute("""
            CREATE TABLE IF NOT EXISTS prereqs (
                course_name TEXT,
                required_course_name TEXT,
                PRIMARY KEY (course_name, required_course_name)
            )
        """)
        # collect known courses and build tolerant maps
        try:
            rows = cur.execute("SELECT course_name, IFNULL(course_code, '') FROM courses").fetchall()
            known = set()
            name_map = {}   # normalized -> actual name
            code_map = {}   # normalized code -> actual name
            for name, code in rows:
                if not name:
                    continue
                known.add(name)
                nclean = name.strip()
                nkey = nclean.lower()
                name_map[nkey] = nclean
                if code:
                    code_map[code.strip().lower()] = nclean
        except Exception:
            # fallback: try schedule table
            try:
                rows = cur.execute("SELECT DISTINCT course_name FROM schedule").fetchall()
                known = {r[0] for r in rows}
                name_map = { (r[0].strip().lower()): r[0] for r in rows if r[0] }
                code_map = {}
            except Exception:
                known = set()
                name_map = {}
                code_map = {}

        # helper to resolve incoming label to an actual known course name if possible
        def resolve_course_label(label):
            if not label:
                return None
            lab = label.strip()
            lnorm = lab.lower()
            # exact match (case sensitive stored name)
            if lab in known:
                return lab
            # normalized name match
            if lnorm in name_map:
                return name_map[lnorm]
            # code match
            if lnorm in code_map:
                return code_map[lnorm]
            # forgiving contains/prefix match against stored names
            for knorm, real in name_map.items():
                if lnorm == knorm or lnorm in knorm or knorm in lnorm:
                    return real
            return None

        for course, req in items:
            try:
                real_course = resolve_course_label(course)
                real_req = resolve_course_label(req)

                if real_course is None or real_req is None:
                    missing_pair = []
                    if real_course is None:
                        missing_pair.append(f"المقرر غير موجود: {course}")
                    if real_req is None:
                        missing_pair.append(f"المقرر المطلوب غير موجود: {req}")
                    missing.append({"course":course,"required":req,"reason":"; ".join(missing_pair)})
                    continue

                if real_course == real_req:
                    errors.append({"course":course,"required":req,"reason":"المقرر لا يمكن أن يكون متطلباً لنفسه"})
                    continue

                cur.execute("INSERT OR IGNORE INTO prereqs (course_name, required_course_name) VALUES (?,?)", (real_course, real_req))
                if cur.rowcount == 0:
                    ignored.append({"course":real_course,"required":real_req})
                else:
                    added.append({"course":real_course,"required":real_req})
            except Exception as e:
                current_app.logger.exception("add_prereq item failed")
                errors.append({"course":course,"required":req,"reason":str(e)})
        conn.commit()

    return jsonify({
        "status":"ok",
        "added": added,
        "ignored": ignored,
        "missing": missing,
        "errors": errors,
        "message": f"تمت المعالجة: تمت الإضافة {len(added)}؛ تجاهل التكرار {len(ignored)}؛ ناقصة {len(missing)}؛ أخطاء {len(errors)}"
    }), 200

@courses_bp.route("/prereqs/delete", methods=["POST"])
@role_required("admin", "admin_main", "head_of_department")
def delete_prereq():
    data = request.get_json(force=True)
    course = data.get("course_name")
    req = data.get("required_course_name")
    if not course or not req:
        return jsonify({"status": "error", "message": "course_name و required_course_name مطلوبة"}), 400
    with get_connection() as conn:
        cur = conn.cursor()
        cur.execute("DELETE FROM prereqs WHERE course_name = ? AND required_course_name = ?", (course, req))
        conn.commit()
    return jsonify({"status": "ok", "message": "تم حذف المتطلب"}), 200

@courses_bp.route("/prereqs/list")
@login_required
def list_prereqs():
    with get_connection() as conn:
        cur = conn.cursor()
        rows = cur.execute("SELECT course_name, required_course_name FROM prereqs ORDER BY course_name, required_course_name").fetchall()
        return jsonify([{"course_name": r[0], "required_course_name": r[1]} for r in rows])

@courses_bp.route("/prereqs/status")
@login_required
def prereq_status():
    student_id = request.args.get("student_id")
    if not student_id:
        return jsonify({"status": "error", "message": "student_id مطلوب"}), 400

    with get_connection() as conn:
        cur = conn.cursor()
        try:
            rows_c = cur.execute("SELECT course_name FROM courses").fetchall()
            courses = [r[0] for r in rows_c]
        except Exception:
            try:
                rows_s = cur.execute("SELECT DISTINCT course_name FROM schedule").fetchall()
                courses = [r[0] for r in rows_s]
            except Exception:
                courses = []

        rows_p = cur.execute("SELECT course_name, required_course_name FROM prereqs").fetchall()
        prereq_map = defaultdict(list)
        for c, req in rows_p:
            prereq_map[c].append(req)

        taken_rows = cur.execute(
            "SELECT DISTINCT course_name FROM grades WHERE student_id = ? AND grade IS NOT NULL", (student_id,)
        ).fetchall()
        taken = {r[0] for r in taken_rows}

        allowed = []
        blocked = {}
        for c in courses:
            reqs = prereq_map.get(c, [])
            missing = [req for req in reqs if req not in taken]
            if missing:
                blocked[c] = missing
            else:
                allowed.append(c)

    return jsonify({"status": "ok", "allowed": allowed, "blocked": blocked, "prereqs": prereq_map})


def _load_courses_and_prereqs(conn):
    cur = conn.cursor()
    try:
        rows_c = cur.execute("SELECT course_name, COALESCE(course_code,'') AS course_code FROM courses").fetchall()
        courses = [{"course_name": r[0], "course_code": (r[1] or "")} for r in (rows_c or []) if r and r[0]]
    except Exception:
        rows_s = cur.execute("SELECT DISTINCT course_name FROM schedule").fetchall()
        courses = [{"course_name": r[0], "course_code": ""} for r in (rows_s or []) if r and r[0]]

    try:
        rows_p = cur.execute("SELECT course_name, required_course_name FROM prereqs").fetchall()
        prereqs = [{"course_name": r[0], "required_course_name": r[1]} for r in (rows_p or []) if r and r[0] and r[1]]
    except Exception:
        prereqs = []

    # normalize + dedupe
    seen = set()
    out_courses = []
    for c in courses:
        key = (c.get("course_name") or "").strip().lower()
        if not key or key in seen:
            continue
        seen.add(key)
        out_courses.append(c)
    seen_p = set()
    out_pr = []
    for p in prereqs:
        a = (p.get("required_course_name") or "").strip()
        b = (p.get("course_name") or "").strip()
        if not a or not b:
            continue
        k = (b.lower(), a.lower())
        if k in seen_p:
            continue
        seen_p.add(k)
        out_pr.append({"course_name": b, "required_course_name": a})
    return out_courses, out_pr


def _subgraph_for_course(prereqs_rows, focus_course: str, direction: str, depth: int):
    focus = (focus_course or "").strip()
    if not focus:
        return prereqs_rows
    direction = (direction or "both").strip().lower()
    if direction not in ("both", "prereqs", "dependents"):
        direction = "both"
    try:
        depth = int(depth or 2)
    except Exception:
        depth = 2
    depth = max(1, min(depth, 10))

    # Build adjacency
    prereq_to_course = defaultdict(set)  # req -> {course}
    course_to_prereq = defaultdict(set)  # course -> {req}
    for row in prereqs_rows:
        c = (row.get("course_name") or "").strip()
        r = (row.get("required_course_name") or "").strip()
        if not c or not r:
            continue
        prereq_to_course[r].add(c)
        course_to_prereq[c].add(r)

    included_courses = {focus}

    def walk_prereqs():
        frontier = {focus}
        for _ in range(depth):
            nxt = set()
            for c in frontier:
                for r in course_to_prereq.get(c, set()):
                    if r not in included_courses:
                        included_courses.add(r)
                        nxt.add(r)
            frontier = nxt
            if not frontier:
                break

    def walk_dependents():
        frontier = {focus}
        for _ in range(depth):
            nxt = set()
            for r in frontier:
                for c in prereq_to_course.get(r, set()):
                    if c not in included_courses:
                        included_courses.add(c)
                        nxt.add(c)
            frontier = nxt
            if not frontier:
                break

    if direction in ("both", "prereqs"):
        walk_prereqs()
    if direction in ("both", "dependents"):
        walk_dependents()

    out = []
    for row in prereqs_rows:
        c = (row.get("course_name") or "").strip()
        r = (row.get("required_course_name") or "").strip()
        if c in included_courses and r in included_courses:
            out.append(row)
    return out


def _render_prereq_flow_png(courses, prereqs_rows, focus_course: str = "", direction: str = "both", depth: int = 2):
    # Import matplotlib lazily so the app runs without it until used.
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception as e:
        raise RuntimeError(
            "ميزة خريطة المتطلبات تحتاج تثبيت مكتبة matplotlib. "
            "نفّذ: pip install -r requirements.txt ثم أعد تشغيل السيرفر."
        ) from e

    # Arabic shaping + bidi so Arabic renders correctly in matplotlib
    try:
        import arabic_reshaper
        from bidi.algorithm import get_display

        def _rtl_text(s: str) -> str:
            txt = str(s or "")
            if not txt:
                return ""
            return get_display(arabic_reshaper.reshape(txt))
    except Exception:
        def _rtl_text(s: str) -> str:
            return str(s or "")

    # Optionally reduce to a focused subgraph
    filtered_pr = _subgraph_for_course(prereqs_rows, focus_course=focus_course, direction=direction, depth=depth)

    # Build nodes set (from edges + courses list)
    nodes = set()
    for p in filtered_pr:
        nodes.add((p.get("course_name") or "").strip())
        nodes.add((p.get("required_course_name") or "").strip())
    nodes = {n for n in nodes if n}
    if not nodes:
        # fallback: show an "empty" image
        fig = plt.figure(figsize=(10, 4), dpi=160)
        ax = fig.add_subplot(111)
        ax.axis("off")
        ax.text(
            0.5, 0.5, _rtl_text("لا توجد متطلبات لعرضها"),
            ha="center", va="center", fontsize=16, fontfamily="DejaVu Sans"
        )
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight")
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()

    # adjacency: prereq -> course
    succ = defaultdict(list)
    indeg = defaultdict(int)
    for n in nodes:
        indeg[n] = 0
    for row in filtered_pr:
        c = (row.get("course_name") or "").strip()
        r = (row.get("required_course_name") or "").strip()
        if not c or not r:
            continue
        if c not in nodes or r not in nodes:
            continue
        succ[r].append(c)
        indeg[c] += 1

    # Layering: Kahn + level propagation (best-effort even with cycles)
    level = {n: 0 for n in nodes}
    q = [n for n in nodes if indeg.get(n, 0) == 0]
    processed = 0
    while q:
        n = q.pop(0)
        processed += 1
        for v in succ.get(n, []):
            level[v] = max(level.get(v, 0), level.get(n, 0) + 1)
            indeg[v] = max(0, indeg.get(v, 0) - 1)
            if indeg[v] == 0:
                q.append(v)
    # cycles: keep existing level=0..N, but still plot

    max_level = max(level.values()) if level else 0
    layers = defaultdict(list)
    for n, lv in level.items():
        layers[lv].append(n)
    for lv in layers:
        layers[lv].sort(key=lambda x: x)

    # Reduce edge crossings: reorder each layer by predecessor barycenter.
    preds = defaultdict(list)  # course -> [prereq]
    for row in filtered_pr:
        c = (row.get("course_name") or "").strip()
        r = (row.get("required_course_name") or "").strip()
        if c and r:
            preds[c].append(r)
    for lv in range(1, max_level + 1):
        prev_order = {n: i for i, n in enumerate(layers.get(lv - 1, []))}
        def _bary(n):
            ps = [prev_order[p] for p in preds.get(n, []) if p in prev_order]
            if not ps:
                return 10**9
            return sum(ps) / max(1, len(ps))
        layers[lv].sort(key=lambda n: (_bary(n), n))

    # coordinates
    pos = {}
    x_scale = 1.45  # more horizontal separation between levels
    for lv in range(0, max_level + 1):
        layer_nodes = layers.get(lv, [])
        for i, n in enumerate(layer_nodes):
            # x increases with level; y decreases with index for top-down
            pos[n] = (lv * x_scale, -i)

    # Figure size scaling
    max_layer_size = max((len(v) for v in layers.values()), default=1)
    width = max(10, 2.5 + (max_level + 1) * 2.8)
    height = max(4.5, 1.6 + max_layer_size * 0.9)

    fig = plt.figure(figsize=(width, height), dpi=160)
    ax = fig.add_subplot(111)
    ax.axis("off")

    # Build display labels with code if available
    code_map = {}
    for c in courses or []:
        name = (c.get("course_name") or "").strip()
        if not name:
            continue
        code_map[name] = (c.get("course_code") or "").strip()

    def label_for(name: str) -> str:
        code = (code_map.get(name) or "").strip()
        raw = f"{name}\n({code})" if code else name
        return _rtl_text(raw)

    # Draw edges first
    is_full_plan = not (focus_course or "").strip()
    edge_palette = ["#0f766e", "#1d4ed8", "#7c3aed", "#b45309", "#be123c", "#0f766e"]
    for row in filtered_pr:
        c = (row.get("course_name") or "").strip()
        r = (row.get("required_course_name") or "").strip()
        if c not in pos or r not in pos:
            continue
        x1, y1 = pos[r]
        x2, y2 = pos[c]
        src_level = int(level.get(r, 0) or 0)
        color = edge_palette[src_level % len(edge_palette)]
        # "core" edges are adjacent levels; others are lighter/dashed.
        level_gap = abs(int(level.get(c, 0) or 0) - src_level)
        core_edge = (level_gap <= 1)
        lw = 1.8 if core_edge else 1.1
        alpha = 0.72 if core_edge else 0.38
        linestyle = "-" if core_edge else "--"
        # In focused mode keep stronger edges for readability.
        if not is_full_plan:
            lw = 1.9 if core_edge else 1.4
            alpha = 0.82 if core_edge else 0.55
            linestyle = "-"
        ax.annotate(
            "",
            xy=(x2, y2),
            xytext=(x1, y1),
            arrowprops=dict(
                arrowstyle="-|>",
                color=color,
                lw=lw,
                linestyle=linestyle,
                alpha=alpha,
                mutation_scale=9,
                shrinkA=12,
                shrinkB=12,
            ),
            zorder=1,
        )

    # Draw nodes
    for n, (x, y) in pos.items():
        is_focus = (focus_course or "").strip() and n == (focus_course or "").strip()
        fc = "#e2e8f0" if not is_focus else "#fde68a"
        ec = "#94a3b8" if not is_focus else "#f59e0b"
        ax.text(
            x, y,
            label_for(n),
            ha="center",
            va="center",
            fontsize=10,
            color="#0f172a",
            bbox=dict(boxstyle="round,pad=0.35", fc=fc, ec=ec, lw=1.2),
            fontfamily="DejaVu Sans",
            zorder=2,
        )

    # Title
    title = "خريطة المتطلبات بين المقررات"
    if (focus_course or "").strip():
        title += f" — ({focus_course})"
    ax.text(
        0, 1.02, _rtl_text(title),
        transform=ax.transAxes,
        ha="left",
        va="bottom",
        fontsize=14,
        fontweight="bold",
        fontfamily="DejaVu Sans",
    )

    # Tight bounds with padding
    xs = [p[0] for p in pos.values()]
    ys = [p[1] for p in pos.values()]
    ax.set_xlim(min(xs) - 1.0, max(xs) + 1.0)
    ax.set_ylim(min(ys) - 1.0, max(ys) + 1.0)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


@courses_bp.route("/prereqs/flowchart/png")
@login_required
def prereqs_flowchart_png():
    course = (request.args.get("course") or "").strip()
    direction = (request.args.get("direction") or "both").strip()
    depth = request.args.get("depth") or 2
    with get_connection() as conn:
        courses, prereqs_rows = _load_courses_and_prereqs(conn)
    try:
        png = _render_prereq_flow_png(courses, prereqs_rows, focus_course=course, direction=direction, depth=depth)
        return send_file(io.BytesIO(png), mimetype="image/png", as_attachment=False, download_name="prereqs_flowchart.png")
    except Exception as e:
        current_app.logger.exception("flowchart png failed")
        return jsonify({"status": "error", "message": str(e)}), 500


@courses_bp.route("/prereqs/flowchart/pdf")
@role_required("admin", "admin_main", "head_of_department", "supervisor", "instructor", "student")
def prereqs_flowchart_pdf():
    course = (request.args.get("course") or "").strip()
    direction = (request.args.get("direction") or "both").strip()
    depth = request.args.get("depth") or 2
    with get_connection() as conn:
        courses, prereqs_rows = _load_courses_and_prereqs(conn)
    try:
        png = _render_prereq_flow_png(courses, prereqs_rows, focus_course=course, direction=direction, depth=depth)
    except Exception as e:
        current_app.logger.exception("flowchart render failed")
        return jsonify({"status": "error", "message": str(e)}), 500
    b64 = base64.b64encode(png).decode("ascii")
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    html = f"""
    <!doctype html>
    <html lang="ar" dir="rtl">
      <head>
        <meta charset="utf-8"/>
        <title>خريطة المتطلبات</title>
        <style>
          body {{ font-family: DejaVu Sans, Arial, Tahoma; direction: rtl; }}
          .meta {{ color:#475569; font-size: 12px; margin-bottom: 8px; }}
          .imgwrap {{ width: 100%; text-align: center; }}
          img {{ max-width: 100%; height: auto; }}
        </style>
      </head>
      <body>
        <h3 style="margin:0 0 6px 0;">خريطة المتطلبات بين المقررات</h3>
        <div class="meta">التاريخ: {now}{(' — مقرر: ' + course) if course else ''}</div>
        <div class="imgwrap">
          <img src="data:image/png;base64,{b64}" alt="Prereqs Flowchart"/>
        </div>
      </body>
    </html>
    """
    return pdf_response_from_html(html, filename_prefix="prereqs_flowchart")


@courses_bp.route("/prereqs/flowchart/pptx")
@role_required("admin", "admin_main", "head_of_department", "supervisor", "instructor", "student")
def prereqs_flowchart_pptx():
    course = (request.args.get("course") or "").strip()
    direction = (request.args.get("direction") or "both").strip()
    depth = request.args.get("depth") or 2
    with get_connection() as conn:
        courses, prereqs_rows = _load_courses_and_prereqs(conn)
    try:
        png = _render_prereq_flow_png(courses, prereqs_rows, focus_course=course, direction=direction, depth=depth)
    except Exception as e:
        current_app.logger.exception("flowchart render failed")
        return jsonify({"status": "error", "message": str(e)}), 500

    # Lazy import so app doesn't crash if dependency missing until used
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as e:
        current_app.logger.exception("pptx dependency missing")
        return jsonify({
            "status": "error",
            "message": "تصدير PowerPoint يحتاج تثبيت python-pptx. نفّذ: pip install -r requirements.txt ثم أعد تشغيل السيرفر."
        }), 500

    prs = Presentation()
    # Use a blank layout if possible
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # Title (textbox)
    title = "خريطة المتطلبات بين المقررات"
    if course:
        title += f" — {course}"
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.6))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True

    # Image
    img_stream = io.BytesIO(png)
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.0), width=Inches(12.5))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return send_file(
        out,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name="prereqs_flowchart.pptx",
    )

# -----------------------
# Export endpoints
# -----------------------

@courses_bp.route("/export/excel")
@login_required
def export_courses_excel():
    """
    Export full courses table as an Excel file using utilities.excel_response_from_df.
    """
    if _is_instructor_or_supervisor_view_only():
        return jsonify({"status": "error", "message": "FORBIDDEN"}), 403
    try:
        df = df_from_query("SELECT course_name, course_code, units FROM courses")
    except Exception:
        # If table doesn't exist or query fails, return empty CSV-like response
        from io import StringIO
        sio = StringIO()
        sio.write("course_name,course_code,units\n")
        sio.seek(0)
        return Response(sio.getvalue(), mimetype="text/csv")
    return excel_response_from_df(df, filename_prefix="courses")

@courses_bp.route("/export/pdf")
@login_required
def export_courses_pdf():
    """
    Export courses list as PDF. If pdf generation is not available, the underlying utility
    will return a JSON error response explaining the issue.
    """
    if _is_instructor_or_supervisor_view_only():
        return jsonify({"status": "error", "message": "FORBIDDEN"}), 403
    try:
        df = df_from_query("SELECT course_name, course_code, units FROM courses")
    except Exception:
        df = None

    if df is None or df.empty:
        html = "<html><head><meta charset='utf-8'><title>المقررات</title></head><body><h3>لا توجد مقررات للتصدير</h3></body></html>"
        return pdf_response_from_html(html, filename_prefix="courses")

    # توليد HTML بسيط من DataFrame (مأمون لمعظم البيانات الصغيرة)
    table_html = df.to_html(index=False, classes="table table-bordered table-sm", border=0, justify="left")
    html = f"""
    <!doctype html>
    <html lang="ar" dir="rtl">
      <head>
        <meta charset="utf-8"/>
        <title>قائمة المقررات</title>
        <style>
          body {{ font-family: DejaVu Sans, Arial, Tahoma; direction: rtl; }}
          table {{ border-collapse: collapse; width: 100%; }}
          table th, table td {{ border: 1px solid #ccc; padding: 6px; text-align: left; }}
          th {{ background: #f0f0f0; }}
        </style>
      </head>
      <body>
        <h3>قائمة المقررات</h3>
        {table_html}
      </body>
    </html>
    """
    return pdf_response_from_html(html, filename_prefix="courses")
